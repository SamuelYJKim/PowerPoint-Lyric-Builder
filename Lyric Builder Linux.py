"""
Samuel Kim
8/14/2020
PowerPoint Lyric Builder
Python file for linux
"""
#!/usr/bin/python
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import lyricsgenius
import PySimpleGUI as sg
import os
import sys

# method for cleaning up string
def remove(string): 
    return string.replace(" ", "") 

# method for converting string into list with \n delimiter
def convert(string): 
    li = list(string.split("\n")) 
    return li

# method for finding how many slides for powerpoint is needed5
def length(lyrics):
	ret = 0
	for line in lyrics:
		if(line == ''):
			ret += 1
	return ret

# method for making an individual slide
def make_slide(prs):
	slide_layout = prs.slide_layouts[6]
	slide = prs.slides.add_slide(slide_layout)
	txBox = slide.shapes.add_textbox(corner, corner, width, height)
	tf = txBox.text_frame
	p = tf.add_paragraph()
	p.alignment = PP_ALIGN.CENTER
	run = p.add_run()
	font = run.font
	font.name = 'Calibri'
	font.size = Pt(48)
	run.font.color.rgb = white_color
	return run

# api key is stored in different file for security reasons
current_dir = os.path.dirname(os.path.realpath(__file__))
api_file = open(current_dir + "/api-key.txt", 'r')
genius = lyricsgenius.Genius(api_file.read())
api_file.close()
genius.remove_section_headers = True
genius.verbose = False
sg.theme('SystemDefault')  

# setting up gui
layout = [
    [sg.Text('Song Title')],
    [sg.Input()],
    [sg.Text('Artist Name')],
    [sg.Input()],
    [sg.Text('Choose Template File')],
    [sg.Input(), sg.FileBrowse( file_types=(('PowerPoint files', '.pptx'),),)],
    [sg.Text('Where do you want to download the file?')],
    [sg.Input(), sg.FolderBrowse()],
    [sg.Text("What is the name of the output file?")],
    [sg.Input()],
    [sg.Text('Log')],
    [sg.Output(size=(45, 5))],
    [sg.Button('Find'), sg.Button('Cancel')],
    ]

# opening gui 
window = sg.Window('PowerPoint Lyric Builder', layout)
while True:
    (event, values) = window.read()
    if event == sg.WIN_CLOSED or event == 'Cancel':  # if user closes window or clicks cancel
        break
    song_name = values[0]
    artist = values[1]
    template_location = values[2]
    file_location = values[3]
    pptx_name = values[4]
    artist_found = genius.search_artist(artist, max_songs=1,
            sort='title')
    if artist_found is None:
        print ('Error, Artist not found!')
        continue
    print("Artist Found")
    song = genius.search_song(song_name, artist_found.name)
    if song is None:
        print ('Error, Song not found!')
        continue
    print("Song Found")
    if template_location == '':
        print("Error, template location not selected")
        continue
    if file_location == '':
        print("Error, file location not selected")
        continue
    name = song.title.replace('(', '')
    name = name.replace(')', '')
    name = name.replace('[', '')
    name = name.replace(']', '')
    lyrics = song.lyrics
    list_lyrics = convert(lyrics)
    current_dir = os.path.dirname(os.path.realpath(__file__))
    
    # Create new presentation
    prs = Presentation(template_location)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    white_color = RGBColor(255, 255, 255)
    title.text = name
    title.text_frame.paragraphs[0].font.color.rgb = white_color
    subtitle.text = artist_found.name
    subtitle.text_frame.paragraphs[0].font.color.rgb = white_color
    
    # Dimensions of the slide.
    corner = Cm(0)
    width = Inches(13.33)
    height = Inches(7.5)
    fileName = remove(name.lower())
    artistFileName = remove(artist_found.name.lower())
    openFile = 'lyrics_' + artistFileName + '_' + fileName + '.txt'
    length_lyrics = length(list_lyrics)
    count = 0
    num = 0
    
    # make each slide and add the lyrics to the paragraph
    while num < length_lyrics:
	    run = make_slide(prs)
	    while 1:
		    line = list_lyrics.pop(0)
		    if(line == ''):
			    break
		    else:
			    run.text += line + '\n'
	    num += 1

    # save powerpoint to directory chosen by user
    prs.save(file_location + '/' + pptx_name + '.pptx')
    print("")
    print("Presentation completed.")
    print("To find new songs, please change entries above.")
    print("Press X or the cancel button to exit")

window.close
